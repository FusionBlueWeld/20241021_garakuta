from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
import datetime

# ==========================================
# 定数設定 (共通スタイル)
# ==========================================
COLOR_MAIN = RGBColor(0, 51, 153)  # 青
COLOR_TEXT = RGBColor(0, 0, 0)     # 黒
COLOR_RED  = RGBColor(255, 0, 0)   # 赤
COLOR_GRAY = RGBColor(128, 128, 128) # 灰
FONT_NAME  = 'BIZ UDPゴシック'

# コンテンツエリアの定義 (ヘッダーとフッターの間)
CONTENT_TOP = Cm(2.5)
CONTENT_LEFT = Cm(2.0)
CONTENT_WIDTH = Cm(29.867) # 33.867 - 4.0
CONTENT_HEIGHT = Cm(15.0)

# ==========================================
# ヘルパー関数
# ==========================================
def set_font_style(paragraph, size, color=COLOR_TEXT, bold=False, font_name=FONT_NAME):
    font = paragraph.font
    font.name = font_name
    font.size = size
    font.bold = bold
    font.color.rgb = color
    paragraph.font._element.set('typeface', font_name)

# ##########################################################################
# 【中身の関数】箇条書きレイアウト
# ##########################################################################
def layout_bullets(slide, subject_text="【主題】ここにテーマを入力", items=None):
    """
    主題(見出し)と箇条書きアイテムを配置する関数
    """
    if items is None:
        items = ["アイテム1", "アイテム2", "アイテム3"]

    # 1. 主題 (青い下線付きの見出し)
    subj_h = Cm(1.2)
    txBox_s = slide.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, subj_h)
    p_s = txBox_s.text_frame.paragraphs[0]
    p_s.text = subject_text
    set_font_style(p_s, size=Pt(20), color=COLOR_MAIN, bold=True)
    
    # 主題の下線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, CONTENT_LEFT, CONTENT_TOP + subj_h, CONTENT_WIDTH, Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_MAIN
    line.line.fill.background()

    # 2. 箇条書き
    bullet_top = CONTENT_TOP + subj_h + Cm(0.5)
    txBox_b = slide.shapes.add_textbox(CONTENT_LEFT + Cm(0.5), bullet_top, CONTENT_WIDTH - Cm(1), CONTENT_HEIGHT - subj_h)
    tf_b = txBox_b.text_frame
    tf_b.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf_b.paragraphs[0]
        else:
            p = tf_b.add_paragraph()
        
        p.text = item
        p.level = 0
        p.space_before = Pt(10)
        set_font_style(p, size=Pt(18))

# ##########################################################################
# 【関数】表紙スライドを作成
# ##########################################################################
def add_cover_slide(prs, title_str="メインタイトル", subtitle_str="サブタイトル", name_str="氏名 太郎"):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 白紙

    # --------------------------------------------------
    # 1. 青い矩形 (左端アクセント)
    # --------------------------------------------------
    obj_w, obj_h = Cm(1), Cm(3)
    obj_left = Cm(2)
    obj_top = (prs.slide_height - obj_h) / 2 
    rect_blue = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, obj_left, obj_top, obj_w, obj_h)
    rect_blue.fill.solid()
    rect_blue.fill.fore_color.rgb = COLOR_MAIN
    rect_blue.line.fill.background()
    rect_blue.shadow.inherit = False

    # --------------------------------------------------
    # 2. ワイドな灰色矩形 (中央ベースライン)
    # --------------------------------------------------
    wide_obj_w, wide_obj_h = Cm(28), Cm(3)
    wide_obj_left = Cm(3.2)
    wide_obj_top = (prs.slide_height - wide_obj_h) / 2
    rect_gray = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, wide_obj_left, wide_obj_top, wide_obj_w, wide_obj_h)
    rect_gray.fill.solid()
    rect_gray.fill.fore_color.rgb = COLOR_GRAY
    rect_gray.line.fill.background()
    rect_gray.shadow.inherit = False

    wide_obj_bottom = wide_obj_top + wide_obj_h

    # --------------------------------------------------
    # 3. メインタイトル (灰色矩形の下端に吸着)
    # --------------------------------------------------
    txBox_t = slide.shapes.add_textbox(wide_obj_left, wide_obj_bottom - Cm(2.5), wide_obj_w, Cm(2.5))
    tf_t = txBox_t.text_frame
    tf_t.vertical_anchor = MSO_ANCHOR.BOTTOM
    p_t = tf_t.paragraphs[0]
    p_t.text = title_str
    set_font_style(p_t, size=Pt(40), color=COLOR_MAIN, bold=True)

    # --------------------------------------------------
    # 4. サブタイトル (灰色矩形の上端に吸着)
    # --------------------------------------------------
    txBox_st = slide.shapes.add_textbox(wide_obj_left, wide_obj_top, wide_obj_w, Cm(1.5))
    tf_st = txBox_st.text_frame
    tf_st.vertical_anchor = MSO_ANCHOR.TOP
    p_st = tf_st.paragraphs[0]
    p_st.text = subtitle_str
    set_font_style(p_st, size=Pt(20), color=COLOR_MAIN, bold=True)

    # --------------------------------------------------
    # 5. 日付 (灰色矩形の1cm下・中央)
    # --------------------------------------------------
    date_str = datetime.datetime.now().strftime("%Y/%m/%d")
    txBox_d = slide.shapes.add_textbox((prs.slide_width - Cm(10))/2, wide_obj_bottom + Cm(1), Cm(10), Cm(1))
    p_d = txBox_d.text_frame.paragraphs[0]
    p_d.text = date_str
    p_d.alignment = PP_ALIGN.CENTER
    set_font_style(p_d, size=Pt(14))

    # --------------------------------------------------
    # 6. 所属・会社名 (灰色矩形の2.5cm下・中央)
    # --------------------------------------------------
    org_text = "パナソニックホールディングス株式会社\n生産技術研究所"
    txBox_o = slide.shapes.add_textbox((prs.slide_width - Cm(20))/2, wide_obj_bottom + Cm(2.5), Cm(20), Cm(2))
    tf_o = txBox_o.text_frame
    tf_o.word_wrap = False
    p_o = tf_o.paragraphs[0]
    p_o.text = org_text
    p_o.alignment = PP_ALIGN.CENTER
    set_font_style(p_o, size=Pt(18))

    # --------------------------------------------------
    # 7. 名前 (灰色矩形の5cm下・中央)
    # --------------------------------------------------
    txBox_n = slide.shapes.add_textbox((prs.slide_width - Cm(10))/2, wide_obj_bottom + Cm(5), Cm(10), Cm(1))
    p_n = txBox_n.text_frame.paragraphs[0]
    p_n.text = name_str
    p_n.alignment = PP_ALIGN.CENTER
    set_font_style(p_n, size=Pt(18))

    return slide


# ##########################################################################
# 【関数】コンテンツスライドを作成
# ##########################################################################
def add_content_slide(prs, title_str="タイトル", page_num_str="#"):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 白紙
    
    header_h = Cm(1.5)
    footer_h = Cm(1.0)

    # --------------------------------------------------
    # 1. ヘッダー：左端の青い矩形
    # --------------------------------------------------
    h_rect_w = Cm(2)
    h_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, h_rect_w, header_h)
    h_rect.fill.solid()
    h_rect.fill.fore_color.rgb = COLOR_MAIN
    h_rect.line.fill.background()
    h_rect.shadow.inherit = False

    # --------------------------------------------------
    # 2. ヘッダー：タイトルテキスト
    # --------------------------------------------------
    t_left, t_width = Cm(2.5), Cm(22)
    txBox_t = slide.shapes.add_textbox(t_left, 0, t_width, header_h)
    tf_t = txBox_t.text_frame
    tf_t.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_t = tf_t.paragraphs[0]
    p_t.text = title_str
    set_font_style(p_t, size=Pt(24), bold=True)

    # --------------------------------------------------
    # 3. ヘッダー：IUO表示 (赤枠)
    # --------------------------------------------------
    iuo_w, iuo_h, iuo_left = Cm(3.35), Cm(1.28), Cm(26)
    iuo_top = (header_h - iuo_h) / 2
    txBox_iuo = slide.shapes.add_textbox(iuo_left, iuo_top, iuo_w, iuo_h)
    txBox_iuo.line.color.rgb, txBox_iuo.line.width = COLOR_RED, Pt(2)
    txBox_iuo.shadow.inherit = False
    tf_iuo = txBox_iuo.text_frame
    tf_iuo.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p1 = tf_iuo.paragraphs[0]
    p1.text = "IUO"
    p1.alignment = PP_ALIGN.CENTER
    set_font_style(p1, size=Pt(16), color=COLOR_RED, bold=True)
    p1.space_after = Pt(0)

    p2 = tf_iuo.add_paragraph()
    p2.text = "Internal Use Only"
    p2.alignment = PP_ALIGN.CENTER
    set_font_style(p2, size=Pt(8), color=COLOR_RED, bold=True)
    p2.space_before = Pt(0)

    # --------------------------------------------------
    # 4. ヘッダー：ページ番号 (#)
    # --------------------------------------------------
    pg_left, pg_width = Cm(32), Cm(1.5)
    txBox_pg = slide.shapes.add_textbox(pg_left, 0, pg_width, header_h)
    tf_pg = txBox_pg.text_frame
    tf_pg.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_pg = tf_pg.paragraphs[0]
    p_pg.text = page_num_str
    set_font_style(p_pg, size=Pt(14))

    # --------------------------------------------------
    # 5. ライン：ヘッダー下の水平線 (全幅)
    # --------------------------------------------------
    line_h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, header_h, prs.slide_width, Pt(2))
    line_h.fill.solid()
    line_h.fill.fore_color.rgb = COLOR_MAIN
    line_h.line.fill.background()
    line_h.shadow.inherit = False

    # --------------------------------------------------
    # 6. ライン：ヘッダー内の縦線 (右から4cm)
    # --------------------------------------------------
    lv_len, lv_x = Cm(1.3), prs.slide_width - Cm(4)
    lv_top = (header_h - lv_len) / 2
    line_v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, lv_x, lv_top, Pt(2), lv_len)
    line_v.fill.solid()
    line_v.fill.fore_color.rgb = COLOR_MAIN
    line_v.line.fill.background()
    line_v.shadow.inherit = False

    # --------------------------------------------------
    # 7. フッター：水平ライン (32cm中央)
    # --------------------------------------------------
    fl_len = Cm(32)
    fl_x = (prs.slide_width - fl_len) / 2
    fl_y = prs.slide_height - footer_h
    line_f = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, fl_x, fl_y, fl_len, Pt(2))
    line_f.fill.solid()
    line_f.fill.fore_color.rgb = COLOR_GRAY
    line_f.line.fill.background()
    line_f.shadow.inherit = False

    # --------------------------------------------------
    # 8. フッター：POWERD BY (左から2cm)
    # --------------------------------------------------
    f_txt_top = prs.slide_height - footer_h
    txBox_f1 = slide.shapes.add_textbox(Cm(2), f_txt_top, Cm(6), footer_h)
    tf_f1 = txBox_f1.text_frame
    tf_f1.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_f1 = tf_f1.paragraphs[0]
    p_f1.text = "POWERD BY\n生産技術研究所"
    set_font_style(p_f1, size=Pt(8))

    # --------------------------------------------------
    # 9. フッター：@2026 Panasonic (左から10cm)
    # --------------------------------------------------
    txBox_f2 = slide.shapes.add_textbox(Cm(10), f_txt_top, Cm(6), footer_h)
    tf_f2 = txBox_f2.text_frame
    tf_f2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_f2 = tf_f2.paragraphs[0]
    p_f2.text = "@2026 Panasonic"
    set_font_style(p_f2, size=Pt(8))

    # --------------------------------------------------
    # 10. フッター：会社表記 (左から22cm)
    # --------------------------------------------------
    txBox_f3 = slide.shapes.add_textbox(Cm(22), f_txt_top, Cm(10), footer_h)
    tf_f3 = txBox_f3.text_frame
    tf_f3.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_f3 = tf_f3.paragraphs[0]
    p_f3.text = "Panasonic Holdings Corporation"
    set_font_style(p_f3, size=Pt(14))

    # --------------------------------------------------
    # 11. コンテンツ描画エリアのガイド（四角のオブジェクト）
    # --------------------------------------------------
    # コンテンツエリア全体を覆う枠線付きの矩形を追加
    guide_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT
    )
    
    # 塗りつぶしを透明（または薄い色）にし、枠線を点線に設定
    guide_rect.fill.background() # 塗りつぶしなし
    guide_rect.line.color.rgb = COLOR_GRAY
    guide_rect.line.width = Pt(1)
    # guide_rect.line.dash_style = 2 # 点線にする場合

    # 範囲内にガイドテキストを表示
    tf_g = guide_rect.text_frame
    tf_g.vertical_anchor = MSO_ANCHOR.MIDDLE
    p_g = tf_g.paragraphs[0]
    p_g.text = "ここにコンテンツを描画するコードを記述してください\n(CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT)"
    p_g.alignment = PP_ALIGN.CENTER
    set_font_style(p_g, size=Pt(16), color=COLOR_GRAY)

    return slide


# ##########################################################################
# メイン実行処理
# ##########################################################################
def main():
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)

    # 1枚目：表紙スライド
    add_cover_slide(
        prs, 
        title_str="自動化設備 開発計画", 
        subtitle_str="～生産効率200%向上を目指して～",
        name_str="氏名 太郎"
    )

    # 2枚目以降：コンテンツスライド
    # スライドを複数追加するテスト
    for i in range(1, 3):
        add_content_slide(
            prs, 
            title_str=f"プロジェクト進捗報告 ({i})", 
            page_num_str=str(i)
        )

    # 保存
    output_file = "Final_Standard_Format.pptx"
    prs.save(output_file)
    print(f"'{output_file}' を作成しました。")

if __name__ == "__main__":
    main()